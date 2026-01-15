import os
import tempfile
import logging
import base64
from pathlib import Path
from typing import List, Dict, Optional, Any, Union
from abc import ABC, abstractmethod
import zipfile
import xml.etree.ElementTree as ET
from io import BytesIO
import json
# Document processing imports
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from docx.document import Document as DocxDoc
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain.text_splitter import MarkdownHeaderTextSplitter
from langchain.schema import Document
import torch
from docling.document_converter import DocumentConverter



logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

if torch.cuda.is_available() and torch.cuda.device_count() > 0:
    logger.info(f"Using GPU: {torch.cuda.get_device_name(0)}")

converter = DocumentConverter()

class BaseDocumentProcessor(ABC):
    """Abstract base class for document processors."""
    
    @abstractmethod
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Process document and return structured content."""
        pass
    
    @abstractmethod
    def extract_images(self, file_path: str, page_number: int = None) -> List[Dict[str, Any]]:
        """Extract images from document."""
        pass

class MultiFormatDocumentExtractor:
    """Enhanced document extractor supporting PDF, DOCX, and PPTX formats."""
    
    SUPPORTED_FORMATS = {'.pdf', '.docx', '.pptx'}
    
    def __init__(self, converter=None):
        """
        Initialize the multi-format extractor.
        
        Args:
            converter: Document converter instance (optional)
        """
        self.converter = converter
        self.file_path = None
        self.processors = {
            '.pdf': PDFProcessor(converter),
            '.docx': DOCXProcessor(),
            '.pptx': PPTXProcessor()
        }
        
    def find_image_indices_with_details(self, text: str, image_marker: str = "<!-- image -->") -> List[Dict[str, Any]]:
        """
        Find all occurrences of image markers with detailed position information.
        
        Args:
            text: The text to search in
            image_marker: The marker to search for
            
        Returns:
            List of dictionaries with occurrence details
            
        Raises:
            ValueError: If text is None or empty
            TypeError: If inputs are not strings
        """
        if not isinstance(text, str):
            raise TypeError("Text must be a string")
        if not isinstance(image_marker, str):
            raise TypeError("Image marker must be a string")
        if not text.strip():
            logger.warning("Empty or whitespace-only text provided")
            return []
        if not image_marker:
            raise ValueError("Image marker cannot be empty")
            
        details = []
        marker_length = len(image_marker)
        occurrence_count = 0
        
        # Pre-compute newline positions for efficient line number calculation
        newline_positions = [-1]  # Start with -1 for column calculation
        for i, char in enumerate(text):
            if char == '\n':
                newline_positions.append(i)
        
        start = 0
        while True:
            try:
                index = text.find(image_marker, start)
                if index == -1:
                    break
                    
                occurrence_count += 1
                end_index = index + marker_length
                
                # Efficient line number calculation using binary search
                line_number = self._get_line_number(newline_positions, index)
                
                # Calculate column position
                line_start = newline_positions[line_number - 1] + 1
                column = index - line_start
                
                details.append({
                    'occurrence': occurrence_count,
                    'start_index': index,
                    'end_index': end_index,
                    'line_number': line_number,
                    'column': column,
                    'marker_text': image_marker
                })
                
                start = end_index
                
            except Exception as e:
                logger.error(f"Error processing occurrence {occurrence_count}: {e}")
                break
                
        logger.info(f"Found {len(details)} image markers in text")
        return details
    
    def _get_line_number(self, newline_positions: List[int], index: int) -> int:
        """Get line number for given character index using binary search."""
        left, right = 0, len(newline_positions) - 1
        
        while left <= right:
            mid = (left + right) // 2
            if newline_positions[mid] < index:
                left = mid + 1
            else:
                right = mid - 1
                
        return left
    
    def load_document(self, file_path: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Load and process document of any supported format.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of processed document pages/slides/sections
            
        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file format is not supported
            RuntimeError: If document processing fails
        """
        if not file_path:
            logger.warning("No file path provided")
            return []
            
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        file_extension = file_path.suffix.lower()
        if file_extension not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported file format: {file_extension}. Supported: {self.SUPPORTED_FORMATS}")
            
        self.file_path = str(file_path)
        
        try:
            processor = self.processors[file_extension]
            result = processor.process_document(self.file_path)
            logger.info(f"Successfully processed {file_extension} document with {len(result)} sections")
            return result
            
        except Exception as e:
            logger.error(f"Failed to process document: {e}")
            raise RuntimeError(f"Document processing failed: {e}")

class PDFProcessor(BaseDocumentProcessor):
    """PDF document processor."""
    
    def __init__(self, converter=None):
        self.converter = converter
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PDF document."""
        try:
            doc = fitz.open(file_path)
            if doc.page_count == 0:
                raise ValueError("PDF document contains no pages")
                
            logger.info(f"Processing PDF with {doc.page_count} pages")
            
            parsed_docs = []
            
            for page_number in range(0, doc.page_count):
                try:
                    page_content = self._process_single_page(doc, page_number, file_path)
                    
                    if page_content:
                        parsed_docs.append(page_content)
                        
                    if page_number % 10 == 0:
                        logger.info(f"Processed {page_number}/{doc.page_count} pages")
                        
                except Exception as e:
                    logger.error(f"Failed to process page {page_number}: {e}")
                    continue
                    
            return parsed_docs
            
        except Exception as e:
            logger.error(f"Failed to process PDF: {e}")
            raise
    
    def write_markdown_to_file(self,file_path, markdown_content):
        """
        Writes markdown content to a file. Appends if the file exists, otherwise creates it.
        
        Args:
            file_path (str): The path to the .md file.
            markdown_content (str): The markdown content to write.
        """
        mode = 'a' if os.path.exists(file_path) else 'w'
        with open(file_path, mode, encoding='utf-8') as f:
            if mode == 'a':
                f.write("\n")  # Optional: add newline before appending
            f.write(markdown_content)

    def _process_single_page(self, doc: fitz.Document, page_number: int, file_path: str) -> Optional[Dict[str, Any]]:
        """Process a single PDF page."""
        try:
            page = doc[page_number]
            
            text_content = page.get_text()

            page_markdown = text_content
            
            temp_page_path = f"Page_{page_number}{file_path}"  

            images = []
            if self.converter:
                    

                    page_document = fitz.open()
                    page_document.insert_pdf(doc, from_page=page_number, to_page=page_number)
                    page_document.save(temp_page_path)
                    page_document.close()

                    conversion_result = converter.convert(temp_page_path)

                    page_markdown = conversion_result.document.export_to_markdown()

                    self.write_markdown_to_file(f"{file_path}.md",page_markdown)

                    if "<!-- image -->" in page_markdown:
                        images = self.extract_images(temp_page_path,page_number)

                        if images:

                            image_details = MultiFormatDocumentExtractor(converter).find_image_indices_with_details(page_markdown)
            
                            page_markdown = self._replace_image_markers(page_markdown, images, image_details)

                            if os.path.exists(temp_page_path):
                                os.remove(temp_page_path)
            
                    
            if os.path.exists(temp_page_path):
                os.remove(temp_page_path)
                            
                        # except Exception as e:
                        #     logger.warning(f"Markdown conversion failed for page {page_number}: {e}")
                        #     page_markdown = text_content
            

            page_content = {
                "text": page_markdown,
                "page_number": page_number,
                "source": file_path,
                "document_type": "pdf",
                "images": images
            }
            
            return page_content

            
        except Exception as e:
            logger.error(f"Error processing PDF page {page_number}: {e}")
            return None
    
    def extract_images(self, file_path:str, page_number:int) -> List[Dict[str, Any]]:
        """Extract images from PDF."""
        try:
            doc = fitz.open(file_path)
        
            images = []
            
            page = doc[0]
            
            page_images = page.get_images(full=False)
            
            for img_index, img in enumerate(page_images):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    
                    if base_image and "image" in base_image:
                        image_info = {
                            "image_name": f"PDF_Page_{page_number}_Image_{img_index}.{base_image.get('ext', 'png')}",
                            "image_type": base_image.get("ext", "png"),
                            "imagebase64": base64.b64encode(base_image["image"]).decode(),
                            "size_bytes": len(base_image["image"]),
                            "image_index": img_index
                        }
                        images.append(image_info)
                        
                except Exception as e:
                    logger.error(f"Failed to extract image {img_index}: {e}")
                    continue
            
            doc.close()
            return images
            
        except Exception as e:
            logger.error(f"Image extraction failed: {e}")
            return []
    
    def _replace_image_markers(self, text: str, images: List[Dict], image_details: List[Dict]) -> str:
        """Replace image markers with image names."""
        modified_text = text
        offset = 0
        
        
        for i, detail in enumerate(image_details):
            if i < len(images):
                start_pos = detail["start_index"] + offset
                end_pos = detail["end_index"] + offset
                image_name = images[i]["image_name"]
                
                if 0 <= start_pos < len(modified_text) and end_pos <= len(modified_text):
                    modified_text = modified_text[:start_pos] + f"[{image_name}]({image_name})" + modified_text[end_pos:]
                    print(f"Replaced marker at {start_pos} with image name: [{image_name}]({image_name})")
                    offset += len(f"[{image_name}]({image_name})") - (end_pos - start_pos)
        return modified_text

class DOCXProcessor(BaseDocumentProcessor):
    """DOCX document processor."""
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Process DOCX document."""
        try:
            doc = DocxDocument(file_path)
            logger.info(f"Processing DOCX document with {len(doc.paragraphs)} paragraphs")
            
            sections = []
            current_section = {
                "text": "",
                "section_number": 1,
                "source": file_path,
                "document_type": "docx",
                "images": [],
                "tables": []
            }
            
            # Process document elements
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    paragraph = Paragraph(element, doc)
                    text = paragraph.text.strip()
                    if text:
                        current_section["text"] += text + "\n"
                        
                elif isinstance(element, CT_Tbl):
                    table = Table(element, doc)
                    table_data = self._extract_table_data(table)
                    current_section["tables"].append(table_data)
            
            # Extract images
            images = self.extract_images(file_path,page_number=None)
            current_section["images"] = images
            
            # Add image markers to text if images exist
            if images:
                for i, img in enumerate(images):
                    current_section["text"] += f"<!-- image -->\n"
            
            sections.append(current_section)
            return sections
            
        except Exception as e:
            logger.error(f"Failed to process DOCX: {e}")
            raise
    
    def extract_images(self, file_path: str, page_number: int = None) -> List[Dict[str, Any]]:
        """Extract images from DOCX document."""
        try:
            images = []
            
            # DOCX files are ZIP archives
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                # Find image files in the media folder
                media_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                
                for i, media_file in enumerate(media_files):
                    try:
                        with docx_zip.open(media_file) as img_file:
                            image_data = img_file.read()
                            
                            # Determine image type from filename
                            img_ext = Path(media_file).suffix[1:] or 'png'
                            
                            image_info = {
                                "image_name": f"DOCX_Image{i}.{img_ext}",
                                "image_type": img_ext,
                                "imagebase64": base64.b64encode(image_data).decode(),
                                "size_bytes": len(image_data),
                                "image_index": i,
                                "source_path": media_file
                            }
                            images.append(image_info)
                            
                    except Exception as e:
                        logger.error(f"Failed to extract image {media_file}: {e}")
                        continue
            
            logger.info(f"Extracted {len(images)} images from DOCX")
            return images
            
        except Exception as e:
            logger.error(f"DOCX image extraction failed: {e}")
            return []
    
    def _extract_table_data(self, table: Table) -> Dict[str, Any]:
        """Extract data from a Word table."""
        table_data = {
            "rows": len(table.rows),
            "columns": len(table.columns) if table.rows else 0,
            "data": []
        }
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data["data"].append(row_data)
        
        return table_data

class PPTXProcessor(BaseDocumentProcessor):
    """PPTX presentation processor."""
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PPTX presentation."""
        try:
            prs = Presentation(file_path)
            logger.info(f"Processing PPTX with {len(prs.slides)} slides")
            
            slides_data = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_content = self._process_slide(slide, slide_number, file_path)
                if slide_content:
                    slides_data.append(slide_content)
            
            return slides_data
            
        except Exception as e:
            logger.error(f"Failed to process PPTX: {e}")
            raise
    
    def _process_slide(self, slide: Slide, slide_number: int, file_path: str) -> Dict[str, Any]:
        """Process a single slide."""
        try:
            slide_content = {
                "text": "",
                "slide_number": slide_number,
                "source": file_path,
                "document_type": "pptx",
                "images": [],
                "shapes": []
            }
            
            # Extract text and shapes
            for shape in slide.shapes:
                shape_info = self._process_shape(shape)
                if shape_info:
                    if shape_info["type"] == "text":
                        slide_content["text"] += shape_info["content"] + "\n"
                    else:
                        slide_content["shapes"].append(shape_info)
            
            # Extract images from this slide
            slide_images = self._extract_slide_images(slide, slide_number)
            slide_content["images"] = slide_images
            
            # Add image markers
            if slide_images:
                for img in slide_images:
                    slide_content["text"] += "<!-- image -->\n"
            
            return slide_content
            
        except Exception as e:
            logger.error(f"Error processing slide {slide_number}: {e}")
            return None
    
    def _process_shape(self, shape: BaseShape) -> Optional[Dict[str, Any]]:
        """Process a shape from the slide."""
        try:
            if hasattr(shape, "text") and shape.text.strip():
                return {
                    "type": "text",
                    "content": shape.text.strip(),
                    "shape_type": str(shape.shape_type)
                }
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return {
                    "type": "image",
                    "shape_type": "picture",
                    "content": "Image placeholder"
                }
            else:
                return {
                    "type": "shape",
                    "shape_type": str(shape.shape_type),
                    "content": f"Shape: {shape.shape_type}"
                }
                
        except Exception as e:
            logger.error(f"Error processing shape: {e}")
            return None
    
    def _extract_slide_images(self, slide: Slide, slide_number: int) -> List[Dict[str, Any]]:
        """Extract images from a specific slide."""
        images = []
        img_index = 0
        
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    
                    image_info = {
                        "image_name": f"PPTX_Slide{slide_number}_Image{img_index}.{image.ext}",
                        "image_type": image.ext,
                        "imagebase64": base64.b64encode(image.blob).decode(),
                        "size_bytes": len(image.blob),
                        "slide_number": slide_number,
                        "image_index": img_index
                    }
                    images.append(image_info)
                    img_index += 1
                    
            except Exception as e:
                logger.error(f"Failed to extract image from slide {slide_number}: {e}")
                continue
        
        return images
    
    def extract_images(self, file_path: str, page_number: int = None) -> List[Dict[str, Any]]:
        """Extract all images from PPTX presentation."""
        try:
            prs = Presentation(file_path)
            all_images = []
            
            slides_to_process = [page_number] if page_number else range(1, len(prs.slides) + 1)
            
            for slide_num in slides_to_process:
                if 1 <= slide_num <= len(prs.slides):
                    slide = prs.slides[slide_num - 1]
                    slide_images = self._extract_slide_images(slide, slide_num)
                    all_images.extend(slide_images)
            
            logger.info(f"Extracted {len(all_images)} images from PPTX")
            return all_images
            
        except Exception as e:
            logger.error(f"PPTX image extraction failed: {e}")
            return []

def save_list_of_dicts_to_json(data, file_path):

    with open(file_path, "w") as f:
        json.dump(data, f, indent=4)


if __name__ == "__main__":
    extractor = MultiFormatDocumentExtractor(converter=converter)
    
    try:

        test_files = [
            "Knowledge-and-Enquiry-Across-Disciplines.pdf",
        ]
        
        for file_path in test_files:
            if os.path.exists(file_path):
                print(f"\n--- Processing {file_path} ---")
                documents = extractor.load_document(file_path)
                print(f"Processed {len(documents)} sections/pages/slides")
                
                # print(documents)
                # Show summary of first section
                if documents:
                    for doc in documents:
                        print(f"Document type: {doc.get('document_type', 'unknown')}")
                        print(f"Text length: {len(doc.get('text', ''))}")
                        print(f"Images found: {len(doc.get('images', []))}")

                save_list_of_dicts_to_json(documents, f"{file_path}.json")
        # Test image marker detection
        test_text = "Header\n<!-- image -->\nMiddle content\n<!-- image -->\nEnd"
        markers = extractor.find_image_indices_with_details(test_text)
        print(f"\nFound {len(markers)} image markers in test text")
        # print(type(documents))
        # print(type(documents[0]))
        # print(documents[0].get('text', 'No text found'))

        # for doc in documents:
        #     print(f"Document type: {doc.get('document_type', 'unknown')}")
        #     print(f"Text: {doc.get('text', '')}")
        #     print(f"Images found: {len(doc.get('images', []))}")
        #     if 'tables' in doc:
        #         print(f"Tables found: {len(doc['tables'])}")
        

        
    except Exception as e:
        logger.error(f"Example execution failed: {e}")